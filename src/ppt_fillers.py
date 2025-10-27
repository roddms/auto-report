import os
import re
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# ---------------------------
# 🔹 Text / Chart Utilities
# ---------------------------

TOKEN_RE = re.compile(r"\{\{([^\{\}]+)\}\}")

def iter_shapes(slide):
    """모든 도형(그룹 내부 포함)을 순회"""
    for shp in slide.shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, "shapes"):
            for s in shp.shapes:
                yield s

# ---------------------------------
# ① 텍스트 토큰 치환
# ---------------------------------
def replace_text_tokens(prs: Presentation, mapping: dict):
    """
    PowerPoint 내 {{TOKEN}} 형태의 텍스트를 mapping 값으로 교체.
    예:
        mapping = {"SL1_VIS_YOY_P_1": "+12.3%", "SL1_SAL_WON_1": "4,070,000,000원"}
    """
    replaced = 0
    for slide in prs.slides:
        for shp in iter_shapes(slide):
            if not getattr(shp, "has_text_frame", False):
                continue
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    before = r.text
                    r.text = TOKEN_RE.sub(
                        lambda m: str(mapping.get(m.group(1).strip(), m.group(0))), r.text
                    )
                    if r.text != before:
                        replaced += 1
    return replaced


# ---------------------------------
# ② 차트 찾기 및 데이터 교체
# ---------------------------------
def find_chart(prs: Presentation, name: str):
    """이름으로 차트 도형 찾기"""
    for slide in prs.slides:
        for shp in iter_shapes(slide):
            if getattr(shp, "name", "") == name and getattr(shp, "has_chart", False):
                return shp.chart
    return None


def replace_chart_data(chart, categories, series_dict):
    """
    차트 데이터 교체 (서식 유지)
    categories : X축 목록
    series_dict : {"시리즈명": [값 리스트], ...}
    """
    cd = ChartData()
    cd.categories = list(categories)
    for sname, svalues in series_dict.items():
        cd.add_series(sname, list(svalues))
    chart.replace_data(cd)


# ---------------------------------
# ③ 표(Table) 교체 (선택사항)
# ---------------------------------
def find_table(prs: Presentation, name: str):
    """이름으로 표 도형 찾기"""
    for slide in prs.slides:
        for shp in iter_shapes(slide):
            if getattr(shp, "name", "") == name and getattr(shp, "has_table", False):
                return shp.table
    return None


def fill_table(table, dataframe):
    """
    pandas DataFrame 데이터를 표에 삽입.
    컬럼 수와 행 수가 PPT 표 구조와 일치해야 함.
    """
    n_rows = len(dataframe)
    n_cols = len(dataframe.columns)
    for r in range(n_rows):
        for c in range(n_cols):
            val = str(dataframe.iat[r, c])
            table.cell(r, c).text = val


# ---------------------------
# 🔹 Formatter Utilities
# ---------------------------

def fmt_int_comma(x):
    """정수 천단위 콤마"""
    try:
        return f"{int(round(float(x))):,}"
    except Exception:
        return str(x)


def fmt_signed_percent_1(x):
    """+/- 표시, 소수점 1자리 %"""
    try:
        return f"{float(x):+0.1f}%"
    except Exception:
        return str(x)


def fmt_won_or_eok(x):
    """1억 이상이면 억원 단위로 변환"""
    try:
        v = float(x)
        if abs(v) >= 1e8:
            return f"{v/1e8:0.1f}억원"
        else:
            return f"{int(v):,}원"
    except Exception:
        return str(x)


FORMATTERS = {
    "int_comma": fmt_int_comma,
    "signed_percent_1": fmt_signed_percent_1,
    "won_or_eok": fmt_won_or_eok,
}


# ---------------------------
#  트리맵, 히트맵 이미지 삽입
# ---------------------------
def add_images_to_presentation(prs: Presentation, image_map: dict):
    """
    image_map: { "도형이름": "이미지경로", ... }
    - 모든 슬라이드를 반복하며, image_map에 있는 도형 이름과 일치하는 
      도형을 찾아 이미지를 삽입합니다.
    """
    import io
    
    # ⭐️ TARGET_SLIDE_INDEX 대신 모든 슬라이드를 반복합니다.
    # prs.slides는 인덱스 0부터 시작합니다.
    for slide_idx, slide in enumerate(prs.slides):
        idx_to_replace = []
        
        # 현재 슬라이드의 모든 도형을 확인
        for i, shp in enumerate(slide.shapes):
            name = shp.name
            # image_map에 현재 도형 이름이 있는지 확인
            if name in image_map:
                # 위치와 크기 정보를 가져옵니다.
                left = shp.left
                top = shp.top
                width = shp.width
                height = shp.height
                idx_to_replace.append((i, name, left, top, width, height))

        # 현재 슬라이드에서 대체할 도형이 있다면 처리
        if not idx_to_replace:
            continue

        print(f"INFO: 슬라이드 {slide_idx+1} ({len(idx_to_replace)}개 차트) 이미지 삽입 시작")
        
        # 뒤에서부터 삭제(인덱스 보존) 및 이미지 삽입
        for i, name, left, top, width, height in reversed(idx_to_replace):
            shp = slide.shapes[i]
            
            # 템플릿 도형(Placeholder) 제거
            slide.shapes._spTree.remove(shp._element)
            
            # 이미지 파일을 바이너리 스트림으로 읽어 삽입 (이전 수정사항 유지)
            img_path = image_map[name]
            
            try:
                with open(img_path, 'rb') as f:
                    image_stream = io.BytesIO(f.read())
                
                # 스트림을 add_picture 함수의 첫 번째 인수로 전달
                slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                
            except FileNotFoundError:
                print(f"⚠️ 경고: 이미지 파일 경로를 찾을 수 없습니다: {img_path}")
            except Exception as e:
                print(f"❌ 이미지 삽입 중 오류 발생 ({img_path}): {e}")


def colorize_arrows(prs):
    """
    PowerPoint 전체 순회하며,
    문단 내에 ▲(상승) 있으면 문단 전체 빨강,
    ▼(하락) 있으면 문단 전체 파랑.
    """
    for slide in prs.slides:
        for shp in iter_shapes(slide):
            if not getattr(shp, "has_text_frame", False):
                continue
            for p in shp.text_frame.paragraphs:
                full_text = "".join(r.text for r in p.runs)
                if "▲" in full_text:
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(231, 76, 60)       # 빨강
                elif "▼" in full_text:
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(0, 112, 192)     # 파랑



def highlight_max_point_only(chart, series_idx=0, max_idx=0,
                             color=RGBColor(231, 76, 60)):
    """
    chart.series[series_idx]의 특정 point(=max_idx)만 강조색 적용.
    나머지 포인트는 건드리지 않아 기존 색 유지됨.
    """
    series = chart.series[series_idx]
    if max_idx < 0 or max_idx >= len(series.points):
        return
    pt = series.points[max_idx]
    fill = pt.format.fill
    fill.solid()
    fill.fore_color.rgb = color

# ---------------------------
# 🔹 Helper (한 번에 실행용)
# ---------------------------
def apply_tokens_and_charts(prs_path, out_path, token_map, chart_map=None, image_map=None):
    """
    PPT에 토큰/차트 모두 반영하고 저장
    token_map: {TOKEN: 값}
    chart_map: {chart_name: (categories, series_dict)}
    """
    prs = Presentation(prs_path)

    replace_text_tokens(prs, token_map)

    colorize_arrows(prs)

    if chart_map:
        for cname, (cats, sdict) in chart_map.items():
            ch = find_chart(prs, cname)
            if ch:
                replace_chart_data(ch, cats, sdict)

                if cname == "SL5_chart_2":
                    try:
                        # 첫 번째 시리즈의 값 리스트 추출
                        vals = next(iter(sdict.values()))
                        if vals:
                            max_idx = vals.index(max(vals))
                            highlight_max_point_only(ch, series_idx=0, max_idx=max_idx, color=RGBColor(231, 76, 60))
                    except Exception as e:
                        print(f"⚠️ SL5_chart_2 강조 처리 실패: {e}")

                if cname == "SL7_chart":
                    try:
                        # 첫 번째 시리즈의 값 리스트 추출
                        vals = next(iter(sdict.values()))
                        if vals:
                            max_idx = vals.index(max(vals))
                            highlight_max_point_only(ch, series_idx=0, max_idx=max_idx, color=RGBColor(91, 155, 213))
                    except Exception as e:
                        print(f"⚠️ SL7_chart 강조 처리 실패: {e}")
    

    if image_map:
        # out_path -> out_path 로 제자리 저장 (같은 경로 덮어씀)
        add_images_to_presentation(prs, image_map)

    prs.save(out_path)
    return out_path